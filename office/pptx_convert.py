import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .images import ImageWriter
from .result import ConvertResult

_MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"


def _table_md(table) -> str:
    rows = list(table.rows)
    if not rows:
        return ""
    lines = []
    header = [c.text.strip() for c in rows[0].cells]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
    return "\n".join(lines)


def _shape_has_math(shape) -> bool:
    if not shape.has_text_frame:
        return False
    el = shape.text_frame._txBody
    return el.find(f".//{{{_MATH_NS}}}oMath") is not None


def _render_shape(shape, title_shape, parts, writer, slide_idx, counters):
    if shape is title_shape:
        return
    stype = shape.shape_type
    if stype == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            _render_shape(sub, title_shape, parts, writer, slide_idx, counters)
        return
    if stype == MSO_SHAPE_TYPE.PICTURE:
        img = shape.image
        counters["img"] += 1
        ref = writer.add(img.blob, img.ext, f"s{slide_idx}_img{counters['img']}")
        parts.append(f"![]({ref})")
        return
    if shape.has_table:
        md = _table_md(shape.table)
        if md:
            parts.append(md)
        return
    if shape.has_text_frame:
        if _shape_has_math(shape):
            parts.append("[equation]")
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            indent = "  " * para.level
            parts.append(f"{indent}- {text}")


def pptx_to_markdown(path: str, *, asset_dir: str | None = None) -> ConvertResult:
    if not os.path.isfile(path):
        raise FileNotFoundError(f"File not found: {path}")
    try:
        prs = Presentation(path)
    except Exception as e:
        raise ValueError(f"Cannot open pptx: {path}") from e

    writer = ImageWriter(asset_dir)
    parts: list = []
    slides = list(prs.slides)
    for idx, slide in enumerate(slides, start=1):
        parts.append(f"<!-- slide {idx} -->")
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.has_text_frame and title_shape.text_frame.text.strip():
            parts.append("# " + title_shape.text_frame.text.strip())
        counters = {"img": 0}
        ordered = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))
        for shape in ordered:
            _render_shape(shape, title_shape, parts, writer, idx, counters)
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append(f"**Notes:** {note}")

    return ConvertResult("\n\n".join(parts), writer.paths, writer.count, len(slides))
