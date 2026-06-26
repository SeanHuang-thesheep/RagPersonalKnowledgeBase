import io

import fitz
import pytest
from pptx import Presentation
from pptx.util import Inches

from office.pptx_convert import pptx_to_markdown


def _png():
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 3, 3))
    pix.set_rect(pix.irect, (10, 20, 30))
    return pix.tobytes("png")


def _build(path):
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    s1.shapes.title.text = "Slide One Title"
    body = s1.placeholders[1].text_frame
    body.text = "first bullet"
    p = body.add_paragraph()
    p.text = "second bullet"
    p.level = 1
    s1.shapes.add_picture(io.BytesIO(_png()), Inches(1), Inches(1))
    s1.notes_slide.notes_text_frame.text = "speaker note here"
    s2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    tbl = s2.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
    tbl.cell(0, 0).text = "h1"
    tbl.cell(0, 1).text = "h2"
    tbl.cell(1, 0).text = "a"
    tbl.cell(1, 1).text = "b"
    prs.save(path)


def test_pptx_basic(tmp_path):
    p = str(tmp_path / "p.pptx")
    _build(p)
    result = pptx_to_markdown(p, asset_dir=str(tmp_path / "assets"))
    md = result.markdown
    assert "<!-- slide 1 -->" in md
    assert "<!-- slide 2 -->" in md
    assert "# Slide One Title" in md
    assert "- first bullet" in md
    assert "second bullet" in md
    assert "![](" in md
    assert "| h1 | h2 |" in md
    assert "**Notes:** speaker note here" in md
    assert result.unit_count == 2
    assert result.image_count == 1


def test_pptx_missing_file():
    with pytest.raises(FileNotFoundError):
        pptx_to_markdown("nope.pptx")
